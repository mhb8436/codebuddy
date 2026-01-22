#!/usr/bin/env python3
"""
CodeBuddy 매뉴얼 PPT 생성 스크립트
- python-pptx를 사용하여 학습자/강사 매뉴얼 PPT 생성
- 스크린샷 이미지를 포함한 프레젠테이션 자동 생성
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 경로 설정
SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
DOCS_DIR = PROJECT_DIR / "docs"
SCREENSHOT_DIR = DOCS_DIR / "screenshots"
OUTPUT_DIR = DOCS_DIR / "manuals"

# 색상 정의
COLOR_PRIMARY = RGBColor(59, 130, 246)      # Blue-500
COLOR_SECONDARY = RGBColor(99, 102, 241)    # Indigo-500
COLOR_DARK = RGBColor(31, 41, 55)           # Gray-800
COLOR_LIGHT = RGBColor(243, 244, 246)       # Gray-100
COLOR_WHITE = RGBColor(255, 255, 255)


# === 학습자 매뉴얼 콘텐츠 ===
LEARNER_SLIDES = [
    {
        "type": "title",
        "title": "CodeBuddy",
        "subtitle": "AI 프로그래밍 튜터\n학습자 매뉴얼",
    },
    {
        "type": "content",
        "title": "목차",
        "bullets": [
            "1. CodeBuddy 소개",
            "2. 회원가입 및 로그인",
            "3. 커리큘럼 선택",
            "4. 개념 학습",
            "5. AI 튜터와 대화하기",
            "6. 연습문제 풀기",
            "7. 시험 보기",
            "8. 코드 실행하기",
        ],
    },
    {
        "type": "section",
        "title": "1. CodeBuddy 소개",
    },
    {
        "type": "content",
        "title": "CodeBuddy란?",
        "bullets": [
            "AI 기반 프로그래밍 학습 플랫폼",
            "JavaScript, TypeScript, Python 학습 지원",
            "수준별 맞춤 학습 (입문/기초/중급)",
            "AI 튜터와 1:1 대화형 학습",
            "실시간 코드 실행 및 피드백",
        ],
    },
    {
        "type": "section",
        "title": "2. 회원가입 및 로그인",
    },
    {
        "type": "image",
        "title": "로그인 화면",
        "image": "learner/01_login.png",
        "description": "이메일과 비밀번호로 로그인합니다.",
    },
    {
        "type": "image",
        "title": "회원가입 화면",
        "image": "learner/02_register.png",
        "description": "강사에게 받은 초대코드를 입력하여 가입합니다.",
    },
    {
        "type": "section",
        "title": "3. 커리큘럼 선택",
    },
    {
        "type": "image",
        "title": "학습 트랙 선택",
        "image": "learner/04_curriculum.png",
        "description": "언어와 수준에 맞는 학습 트랙을 선택합니다.",
    },
    {
        "type": "section",
        "title": "4. 개념 학습",
    },
    {
        "type": "image",
        "title": "개념 학습 화면",
        "image": "learner/05_concept.png",
        "description": "각 토픽의 핵심 개념과 예제 코드를 학습합니다.",
    },
    {
        "type": "section",
        "title": "5. AI 튜터",
    },
    {
        "type": "image",
        "title": "AI 튜터와 대화하기",
        "image": "learner/06_ai_tutor.png",
        "description": "궁금한 점을 AI에게 질문하세요.\n코드 에디터의 코드를 참조하여 설명해줍니다.",
    },
    {
        "type": "content",
        "title": "AI 튜터 활용 팁",
        "bullets": [
            '"이 코드 설명해줘" - 에디터 코드 설명',
            '"for문이 뭐야?" - 개념 질문',
            '"이 코드 개선해줘" - 코드 리뷰',
            '"에러가 나는데 왜 그래?" - 디버깅 도움',
        ],
    },
    {
        "type": "section",
        "title": "6. 연습문제",
    },
    {
        "type": "image",
        "title": "연습문제 풀기",
        "image": "learner/07_practice.png",
        "description": "AI가 생성한 맞춤형 연습문제를 풀어봅니다.",
    },
    {
        "type": "section",
        "title": "7. 시험",
    },
    {
        "type": "image",
        "title": "시험 보기",
        "image": "learner/08_exam.png",
        "description": "학습한 내용을 시험으로 점검합니다.",
    },
    {
        "type": "content",
        "title": "학습 완료!",
        "bullets": [
            "개념 학습 -> AI 튜터 질문 -> 연습문제 -> 시험",
            "막히는 부분은 언제든 AI에게 질문하세요",
            "코드를 직접 실행하며 익히는 것이 중요합니다",
            "Happy Coding!",
        ],
    },
]


# === 강사 매뉴얼 콘텐츠 ===
INSTRUCTOR_SLIDES = [
    {
        "type": "title",
        "title": "CodeBuddy",
        "subtitle": "AI 프로그래밍 튜터\n강사 매뉴얼",
    },
    {
        "type": "content",
        "title": "목차",
        "bullets": [
            "1. 강사 기능 개요",
            "2. 관리자 접속",
            "3. 반 생성 및 관리",
            "4. 학생 초대 (초대코드)",
            "5. 학생 진도 확인",
            "6. FAQ",
        ],
    },
    {
        "type": "section",
        "title": "1. 강사 기능 개요",
    },
    {
        "type": "content",
        "title": "강사가 할 수 있는 것",
        "bullets": [
            "반(Class) 생성 및 관리",
            "학생 초대코드 발급",
            "학생 목록 및 진도 확인",
            "커리큘럼 설정",
            "AI 모델 설정 (선택)",
        ],
    },
    {
        "type": "section",
        "title": "2. 관리자 접속",
    },
    {
        "type": "image",
        "title": "관리자 대시보드",
        "image": "instructor/admin_01_dashboard.png",
        "description": "상단 '관리자' 메뉴를 클릭하여 접속합니다.",
    },
    {
        "type": "section",
        "title": "3. 반 생성 및 관리",
    },
    {
        "type": "image",
        "title": "반 관리 화면",
        "image": "instructor/admin_02_classes.png",
        "description": "새 반을 생성하고 초대코드를 확인합니다.",
    },
    {
        "type": "content",
        "title": "반 생성 방법",
        "bullets": [
            "1. '반 관리' 메뉴 클릭",
            "2. '새 반 만들기' 버튼 클릭",
            "3. 반 이름 입력 (예: '웹개발 심화반')",
            "4. 최대 인원 설정",
            "5. 생성 완료 후 초대코드 확인",
        ],
    },
    {
        "type": "section",
        "title": "4. 학생 초대",
    },
    {
        "type": "content",
        "title": "초대코드 안내",
        "bullets": [
            "반 생성 시 6자리 초대코드 자동 생성",
            "예: K7X2P9",
            "학생에게 오프라인으로 전달",
            "학생은 회원가입 시 초대코드 입력",
            "코드 노출 시 '코드 재생성' 가능",
        ],
    },
    {
        "type": "section",
        "title": "5. 진도 확인",
    },
    {
        "type": "content",
        "title": "학생 진도 확인 방법",
        "bullets": [
            "반 상세 페이지에서 학생 목록 확인",
            "각 학생의 학습 진행률 표시",
            "완료한 토픽 / 전체 토픽 비율",
            "연습문제 정답률 확인 가능",
        ],
    },
    {
        "type": "section",
        "title": "6. FAQ",
    },
    {
        "type": "content",
        "title": "자주 묻는 질문",
        "bullets": [
            "Q: 초대코드가 노출되었어요",
            "A: 반 관리 > 코드 재생성 클릭",
            "",
            "Q: 학생이 다른 반으로 이동해야 해요",
            "A: 현재는 직접 DB 수정 필요 (추후 기능 추가 예정)",
            "",
            "Q: AI 모델을 변경하고 싶어요",
            "A: 관리자 > 모델 설정에서 수준별 모델 변경 가능",
        ],
    },
]


def create_title_slide(prs: Presentation, title: str, subtitle: str):
    """타이틀 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # 배경색 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER


def create_section_slide(prs: Presentation, title: str):
    """섹션 구분 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # 배경색 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_SECONDARY

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER


def create_content_slide(prs: Presentation, title: str, bullets: list):
    """내용 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    # 불릿 포인트
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet.startswith("Q:") or bullet.startswith("A:"):
            p.text = bullet
            p.font.bold = bullet.startswith("Q:")
        else:
            p.text = f"• {bullet}" if bullet else ""

        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(12)


def create_image_slide(prs: Presentation, title: str, image_path: str, description: str):
    """이미지 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    # 이미지
    full_image_path = SCREENSHOT_DIR / image_path
    if full_image_path.exists():
        # 이미지 크기 계산 (가로 기준 맞춤)
        slide.shapes.add_picture(
            str(full_image_path),
            Inches(0.5), Inches(1),
            width=Inches(9)
        )
    else:
        # 이미지 없을 때 플레이스홀더
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1),
            Inches(9), Inches(5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_LIGHT
        shape.line.color.rgb = COLOR_DARK

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[스크린샷: {image_path}]"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_DARK
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

    # 설명
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_DARK
    p.alignment = PP_ALIGN.CENTER


def generate_ppt(slides_config: list, output_filename: str):
    """PPT 파일 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_config in slides_config:
        slide_type = slide_config["type"]

        if slide_type == "title":
            create_title_slide(prs, slide_config["title"], slide_config["subtitle"])
        elif slide_type == "section":
            create_section_slide(prs, slide_config["title"])
        elif slide_type == "content":
            create_content_slide(prs, slide_config["title"], slide_config["bullets"])
        elif slide_type == "image":
            create_image_slide(
                prs,
                slide_config["title"],
                slide_config["image"],
                slide_config["description"]
            )

    output_path = OUTPUT_DIR / output_filename
    prs.save(str(output_path))
    print(f"✓ PPT 생성 완료: {output_path}")
    return output_path


def main():
    """메인 실행 함수"""
    # 출력 디렉토리 생성
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    print("=" * 50)
    print("CodeBuddy 매뉴얼 PPT 생성")
    print("=" * 50)

    # 학습자 매뉴얼 생성
    print("\n[학습자 매뉴얼 생성 중...]")
    learner_ppt = generate_ppt(LEARNER_SLIDES, "CodeBuddy_학습자_매뉴얼.pptx")

    # 강사 매뉴얼 생성
    print("\n[강사 매뉴얼 생성 중...]")
    instructor_ppt = generate_ppt(INSTRUCTOR_SLIDES, "CodeBuddy_강사_매뉴얼.pptx")

    print("\n" + "=" * 50)
    print("PPT 생성 완료!")
    print(f"학습자 매뉴얼: {learner_ppt}")
    print(f"강사 매뉴얼: {instructor_ppt}")
    print("=" * 50)


if __name__ == "__main__":
    main()
